# file_extract_router.py
from __future__ import annotations

import base64
import json
import os
import shutil
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup

try:
    from dotenv import load_dotenv

    load_dotenv()
except Exception:
    pass

TEXT_EXTS = {".txt", ".md", ".log"}
HTML_EXTS = {".html", ".htm"}
JSON_EXTS = {".json"}
CSV_EXTS = {".csv", ".tsv"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}
DOCX_EXTS = {".docx"}
WORD_EXTS = {".doc", ".docx"}
PPTX_EXTS = {".pptx"}
PPT_EXTS = {".ppt", ".pptx"}
XLSX_EXTS = {".xlsx", ".xls"}
PDF_EXTS = {".pdf"}
OFFICE_EXTS = {".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".odt", ".ods", ".odp"}

# OCR triggers only when PDF looks scanned/image-heavy.
MIN_TEXT_CHARS = 20
IMAGE_PAGE_RATIO = 0.6


def extract_lines_any(file_path: str | Path, ocr_lang: str = "eng") -> list[dict]:
    """
    Returns list of {"page": int, "text": str}
    - PDF: real page numbers
    - Other formats: page=1
    """
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")

    ext = p.suffix.lower()

    # If PDF, keep per-page lines
    if ext in PDF_EXTS:
        return _extract_pdf_lines(p, ocr_lang=ocr_lang)

    if ext in WORD_EXTS:
        with tempfile.TemporaryDirectory(prefix="doc_to_pdf_") as tmp:
            pdf_path = _convert_to_pdf(p, Path(tmp))
            lines = _extract_pdf_lines(pdf_path, ocr_lang=ocr_lang)
            _write_ocr_md(p, "\n".join(d["text"] for d in lines))
            return lines
    if ext in PPT_EXTS:
        with tempfile.TemporaryDirectory(prefix="ppt_to_pdf_") as tmp:
            pdf_path = _convert_to_pdf(p, Path(tmp))
            lines = _extract_pdf_lines(pdf_path, ocr_lang=ocr_lang)
            _write_ocr_md(p, "\n".join(d["text"] for d in lines))
            return lines

    # Non-PDF: extract full text and treat as page 1
    text = extract_text(p, ocr_lang=ocr_lang)
    return [{"page": 1, "text": ln} for ln in _nonempty_lines(text)]


def extract_pdf_pages(file_path: str | Path, ocr_lang: str = "eng") -> list[dict]:
    """
    Returns list of {"page": int, "text": str} for PDFs.
    Uses OCR when needed (or when OCR_ALWAYS=1).
    """
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")
    if p.suffix.lower() not in PDF_EXTS:
        raise ValueError(f"Not a PDF: {p}")

    doc = fitz.open(str(p))
    pages_text = []
    image_pages = 0
    for i in range(len(doc)):
        page = doc[i]
        t = page.get_text("text") or ""
        pages_text.append(t)
        if page.get_images(full=True):
            image_pages += 1
    doc.close()

    ocr_always = os.environ.get("OCR_ALWAYS", "0").strip().lower() in {"1", "true", "yes"}
    if ocr_always or _pdf_needs_ocr(pages_text, image_pages, total_pages=len(pages_text)):
        pages = _build_ocr_pages(p, lang=ocr_lang, dpi=int(os.environ.get("OCR_DPI", "220")))
        if pages:
            return _ocr_pages_to_text(pages)

    # Fallback: real text per page
    out = []
    for page_no, t in enumerate(pages_text, start=1):
        out.append({"page": page_no, "text": t})
    return out


def extract_text(file_path: str | Path, ocr_lang: str = "eng") -> str:
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")

    ext = p.suffix.lower()

    if ext in PDF_EXTS:
        return _extract_pdf(p, ocr_lang=ocr_lang)
    if ext in WORD_EXTS:
        with tempfile.TemporaryDirectory(prefix="doc_to_pdf_") as tmp:
            pdf_path = _convert_to_pdf(p, Path(tmp))
            text = _extract_pdf(pdf_path, ocr_lang=ocr_lang)
            _write_ocr_md(p, text)
            return text
    if ext in PPT_EXTS:
        with tempfile.TemporaryDirectory(prefix="ppt_to_pdf_") as tmp:
            pdf_path = _convert_to_pdf(p, Path(tmp))
            text = _extract_pdf(pdf_path, ocr_lang=ocr_lang)
            _write_ocr_md(p, text)
            return text
    if ext in DOCX_EXTS:
        text = _extract_docx(p)
        _write_ocr_md(p, text)
        return text
    if ext in XLSX_EXTS:
        text = _extract_spreadsheet(p)
        _write_ocr_md(p, text)
        return text
    if ext in CSV_EXTS:
        text = _extract_csv(p)
        _write_ocr_md(p, text)
        return text
    if ext in JSON_EXTS:
        text = _extract_json(p)
        _write_ocr_md(p, text)
        return text
    if ext in HTML_EXTS:
        text = _extract_html(p)
        _write_ocr_md(p, text)
        return text
    if ext in TEXT_EXTS:
        text = p.read_text(encoding="utf-8", errors="ignore")
        _write_ocr_md(p, text)
        return text
    if ext in IMAGE_EXTS:
        return _ocr_image(p, lang=ocr_lang)

    text = p.read_text(encoding="utf-8", errors="ignore")
    _write_ocr_md(p, text)
    return text


def extract_flowchart_graph(
    file_path: str | Path, *, dpi: int = 300, ocr_lang: str = "eng", engine: str | None = None
) -> dict:
    """
    Returns {"nodes": [...], "edges": [...], "pages": int} for PDF/image inputs.
    engine: "groq" or "local". Defaults to "groq" if GROQ_API_KEY is set.
    """
    if engine is None:
        engine = "groq" if os.environ.get("GROQ_API_KEY") else "local"

    if engine == "groq":
        from groq_flowchart import extract_flowchart_graph as _extract_graph
        return _extract_graph(file_path, dpi=dpi)

    from flowchart_graph import extract_flowchart_graph as _extract_graph

    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")
    return _extract_graph(p, dpi=dpi, ocr_lang=ocr_lang)


def ensure_pdf_path(file_path: str | Path, out_dir: Path | None = None, *, ocr_lang: str = "eng") -> Path:
    """
    Convert any supported input to a PDF and return its path.
    - PDFs are returned as-is.
    - Office files use LibreOffice conversion.
    - Images are embedded into a PDF page.
    - Other formats are converted to text and written into a PDF.
    """
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")

    if p.suffix.lower() in PDF_EXTS:
        return p

    if out_dir is None:
        out_dir = Path(tempfile.mkdtemp(prefix="to_pdf_"))
    out_dir.mkdir(parents=True, exist_ok=True)

    ext = p.suffix.lower()
    if ext in OFFICE_EXTS:
        return _convert_to_pdf(p, out_dir)
    if ext in IMAGE_EXTS:
        return _image_to_pdf(p, out_dir)

    text = extract_text(p, ocr_lang=ocr_lang)
    return _text_to_pdf(text, p, out_dir)


# -------------------------
# Extractors
# -------------------------

def _extract_pdf(p: Path, ocr_lang: str = "eng") -> str:
    doc = fitz.open(str(p))
    pages_text = []
    image_pages = 0

    for i in range(len(doc)):
        page = doc[i]
        t = page.get_text("text") or ""
        pages_text.append(t)
        if page.get_images(full=True):
            image_pages += 1
    doc.close()

    if _pdf_needs_ocr(pages_text, image_pages, total_pages=len(pages_text)):
        return _ocr_pdf_pages(p, lang=ocr_lang, dpi=220)

    text = "\n\n".join(pages_text)
    _write_ocr_md(p, text)
    return text


def _extract_pdf_lines(p: Path, ocr_lang: str = "eng") -> list[dict]:
    doc = fitz.open(str(p))
    out = []
    pages_text = []
    image_pages = 0

    for i in range(len(doc)):
        page = doc[i]
        t = page.get_text("text") or ""
        pages_text.append(t)
        if page.get_images(full=True):
            image_pages += 1
    doc.close()

    if _pdf_needs_ocr(pages_text, image_pages, total_pages=len(pages_text)):
        text = _ocr_pdf_pages(p, lang=ocr_lang, dpi=220)
        return [{"page": 1, "text": ln} for ln in _nonempty_lines(text)]

    _write_ocr_md(p, "\n\n".join(pages_text))
    for page_no, t in enumerate(pages_text, start=1):
        for ln in _nonempty_lines(t):
            out.append({"page": page_no, "text": ln})
    return out


def _ocr_pdf_pages(p: Path, lang: str = "eng", dpi: int = 220) -> str:
    _ = lang
    try:
        dpi = int(os.environ.get("OCR_DPI", str(dpi)))
    except Exception:
        pass
    graph_text = _maybe_groq_flowchart_text(p, dpi=dpi)
    if graph_text is not None:
        _write_searchable_pdf(p, lang=lang, dpi=dpi)
        pages = _build_ocr_pages(p, lang=lang, dpi=dpi)
        _write_positioned_html(p, pages=pages)
        layout_text = _write_layout_md(p, pages=pages)
        _append_layout_to_md(p, layout_text)
        return graph_text
    text = _docling_rapidocr_markdown(p)
    _write_ocr_md(p, text)
    _write_searchable_pdf(p, lang=lang, dpi=dpi)
    pages = _build_ocr_pages(p, lang=lang, dpi=dpi)
    _write_positioned_html(p, pages=pages)
    layout_text = _write_layout_md(p, pages=pages)
    _append_layout_to_md(p, layout_text)
    print(f"[flowchart] wrote {p.with_suffix('.ocr.md')}")
    return text


def _extract_docx(p: Path) -> str:
    doc = Document(str(p))
    parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    return "\n".join(parts)


def _extract_pptx(p: Path) -> str:
    prs = Presentation(str(p))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_spreadsheet(p: Path) -> str:
    parts = []
    xls = pd.ExcelFile(str(p))
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name=sheet_name)
        parts.append(f"--- Sheet: {sheet_name} ---")
        parts.append(df.to_string(index=False))
    return "\n\n".join(parts)


def _extract_csv(p: Path) -> str:
    sep = "\t" if p.suffix.lower() == ".tsv" else ","
    df = pd.read_csv(str(p), sep=sep, dtype=str, keep_default_na=False)
    return df.to_string(index=False)


def _extract_json(p: Path) -> str:
    obj = json.loads(p.read_text(encoding="utf-8", errors="ignore"))
    return json.dumps(obj, ensure_ascii=False, indent=2)


def _extract_html(p: Path) -> str:
    html = p.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def _ocr_image(p: Path, lang: str = "eng") -> str:
    _ = lang
    try:
        dpi = int(os.environ.get("OCR_DPI", "220"))
    except Exception:
        dpi = 220
    graph_text = _maybe_groq_flowchart_text(p, dpi=dpi)
    if graph_text is not None:
        _write_searchable_pdf(p, lang=lang, dpi=dpi)
        pages = _build_ocr_pages(p, lang=lang, dpi=dpi)
        _write_positioned_html(p, pages=pages)
        layout_text = _write_layout_md(p, pages=pages)
        _append_layout_to_md(p, layout_text)
        return graph_text
    text = _docling_rapidocr_markdown(p)
    _write_ocr_md(p, text)
    _write_searchable_pdf(p, lang=lang, dpi=dpi)
    pages = _build_ocr_pages(p, lang=lang, dpi=dpi)
    _write_positioned_html(p, pages=pages)
    layout_text = _write_layout_md(p, pages=pages)
    _append_layout_to_md(p, layout_text)
    print(f"[flowchart] wrote {p.with_suffix('.ocr.md')}")
    return text


def _nonempty_lines(text: str):
    for ln in text.splitlines():
        ln = ln.strip()
        if ln:
            yield ln


def _pdf_needs_ocr(pages_text: list[str], image_pages: int, total_pages: int) -> bool:
    total_chars = sum(len(t.strip()) for t in pages_text)
    text_pages = sum(1 for t in pages_text if t.strip())

    if total_pages == 0:
        return False
    if image_pages == 0:
        return False
    if text_pages == 0:
        return True
    if total_chars < MIN_TEXT_CHARS and image_pages >= max(1, int(total_pages * IMAGE_PAGE_RATIO)):
        return True
    return False


def _ocr_pages_to_text(
    pages: list[
        tuple["Image.Image", list[tuple[int, int, int, int, str, int, int, int]], float]
    ]
) -> list[dict]:
    out = []
    for page_no, (_img, words, _avg_conf) in enumerate(pages, start=1):
        words_by_line: dict[tuple[int, int, int], list[tuple[int, str]]] = {}
        for x1, _y1, _x2, _y2, text, block, par, line in words:
            text = (text or "").strip()
            if not text:
                continue
            key = (block, par, line)
            words_by_line.setdefault(key, []).append((x1, text))
        lines = []
        for key in sorted(words_by_line.keys()):
            parts = [t for _x, t in sorted(words_by_line[key], key=lambda v: v[0])]
            if parts:
                lines.append(" ".join(parts))
        out.append({"page": page_no, "text": "\n".join(lines)})
    return out


def _docling_rapidocr_markdown(p: Path) -> str:
    from docling_to_md_auto_rapidocr import convert_to_markdown

    return convert_to_markdown(p)


def _maybe_groq_flowchart_text(p: Path, *, dpi: int) -> str | None:
    if not os.environ.get("GROQ_API_KEY"):
        print("[flowchart] GROQ_API_KEY not set; skipping Groq flowchart")
        return None
    if os.environ.get("FORCE_GROQ_FLOWCHART") == "1":
        print("[flowchart] FORCE_GROQ_FLOWCHART=1; bypassing heuristic")
    elif not _looks_like_flowchart(p, dpi=dpi):
        print("[flowchart] heuristic negative; skipping Groq flowchart")
        return None
    try:
        from groq_flowchart import extract_flowchart_graph as _extract_groq_graph
        from groq_flowchart import to_mermaid_markdown as _to_mermaid
    except Exception:
        print("[flowchart] failed to import groq_flowchart; skipping")
        return None
    try:
        graph = _extract_groq_graph(p, dpi=dpi)
    except Exception:
        print("[flowchart] Groq call failed; skipping")
        return None
    if not graph.get("nodes") and not graph.get("edges"):
        print("[flowchart] Groq returned empty graph; falling back to OCR")
        ocr_text = _docling_rapidocr_markdown(p)
        _write_ocr_md(p, ocr_text)
        return ocr_text
    mermaid = _to_mermaid(graph)
    ocr_text = _docling_rapidocr_markdown(p)
    md = (
        "```mermaid\n"
        + "\n".join(mermaid.strip().splitlines()[1:-1])
        + "\n```\n\n"
        "```text\n"
        + (ocr_text.strip() + "\n" if ocr_text.strip() else "")
        + "```\n"
    )
    _write_flowchart_md(p, md)
    print(f"[flowchart] wrote {p.with_suffix('.flowchart.md')}")
    return ocr_text


def _looks_like_flowchart(p: Path, *, dpi: int) -> bool:
    try:
        import cv2
        import numpy as np
    except Exception:
        return False

    img = _render_first_page_or_image(p, dpi=dpi)
    if img is None:
        return False

    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    blur = cv2.GaussianBlur(gray, (3, 3), 0)
    edges = cv2.Canny(blur, 50, 150)
    _, thresh = cv2.threshold(blur, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

    contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    h, w = edges.shape[:2]
    min_area = max(150, int(0.001 * h * w))

    rect_like = 0
    arrow_like = 0
    for c in contours:
        area = cv2.contourArea(c)
        if area < min_area:
            continue
        peri = cv2.arcLength(c, True)
        approx = cv2.approxPolyDP(c, 0.02 * peri, True)
        if 4 <= len(approx) <= 6:
            rect_like += 1
        elif len(approx) == 3:
            # Triangles often indicate arrowheads.
            if area < 0.02 * h * w:
                arrow_like += 1

    lines = cv2.HoughLinesP(edges, 1, np.pi / 180, threshold=40, minLineLength=30, maxLineGap=8)
    line_count = 0 if lines is None else len(lines)
    hv_lines = 0
    diag_lines = 0
    if lines is not None:
        for x1, y1, x2, y2 in lines[:, 0]:
            dx = x2 - x1
            dy = y2 - y1
            if dx == 0 and dy == 0:
                continue
            angle = abs(np.degrees(np.arctan2(dy, dx)))
            # Normalize to 0-90
            angle = angle if angle <= 90 else 180 - angle
            if angle <= 10 or angle >= 80:
                hv_lines += 1
            else:
                diag_lines += 1

    # Heuristic: flowcharts tend to have multiple boxes + connector lines.
    print(
        "[flowchart] heuristic rects="
        f"{rect_like} arrows={arrow_like} lines={line_count} hv={hv_lines} "
        f"diag={diag_lines} (dpi={dpi})"
    )

    # Table-like grids: lots of horizontal/vertical lines and few diagonals.
    if hv_lines >= 40 and diag_lines == 0 and arrow_like == 0 and rect_like < 3:
        return False

    if arrow_like >= 1:
        return True
    return rect_like >= 1 or (diag_lines >= 1 or line_count >= 40)


def _render_first_page_or_image(p: Path, *, dpi: int):
    suffix = p.suffix.lower()
    try:
        if suffix in PDF_EXTS:
            doc = fitz.open(str(p))
            if len(doc) == 0:
                doc.close()
                return None
            page = doc[0]
            zoom = dpi / 72.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            doc.close()
            img = pix.tobytes("png")
            import numpy as np
            import cv2
            data = np.frombuffer(img, dtype=np.uint8)
            return cv2.imdecode(data, cv2.IMREAD_COLOR)
        if suffix in IMAGE_EXTS:
            import cv2
            return cv2.imread(str(p))
    except Exception:
        return None
    return None


def _image_to_pdf(p: Path, out_dir: Path) -> Path:
    out_path = out_dir / f"{p.stem}.pdf"
    pix = fitz.Pixmap(str(p))
    doc = fitz.open()
    try:
        page = doc.new_page(width=pix.width, height=pix.height)
        page.insert_image(page.rect, filename=str(p))
        doc.save(out_path, deflate=True)
    finally:
        doc.close()
    return out_path


def _text_to_pdf(text: str, src: Path, out_dir: Path) -> Path:
    out_path = out_dir / f"{src.stem}.pdf"
    doc = fitz.open()
    try:
        page_w, page_h = 612, 792  # Letter
        margin = 40
        font_size = 11
        line_height = 14
        max_lines = max(1, int((page_h - 2 * margin) / line_height))
        lines = text.splitlines() or [""]
        for i in range(0, len(lines), max_lines):
            chunk = "\n".join(lines[i:i + max_lines])
            page = doc.new_page(width=page_w, height=page_h)
            rect = fitz.Rect(margin, margin, page_w - margin, page_h - margin)
            page.insert_textbox(rect, chunk, fontsize=font_size, fontname="helv")
        doc.save(out_path, deflate=True)
    finally:
        doc.close()
    return out_path


def _write_flowchart_md(p: Path, content: str) -> None:
    out_path = p.with_suffix(".flowchart.md")
    out_path.write_text(content, encoding="utf-8")


def _write_ocr_md(p: Path, content: str) -> None:
    out_path = p.with_suffix(".ocr.md")
    out_path.write_text(content, encoding="utf-8")
    print(f"[flowchart] wrote {out_path}")


def _write_searchable_pdf(p: Path, *, lang: str, dpi: int) -> Path | None:
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        print("[ocr] pytesseract/Pillow not available; skipping searchable PDF")
        return None

    out_path = p.with_suffix(".ocr.pdf")
    try:
        if p.suffix.lower() in PDF_EXTS:
            src = fitz.open(str(p))
            out_pdf = fitz.open()
            try:
                zoom = dpi / 72.0
                mat = fitz.Matrix(zoom, zoom)
                for page in src:
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    img = Image.open(BytesIO(pix.tobytes("png")))
                    pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                        img, extension="pdf", lang=lang
                    )
                    ocr_page = fitz.open("pdf", pdf_bytes)
                    out_pdf.insert_pdf(ocr_page)
                    ocr_page.close()
                out_pdf.save(out_path, deflate=True)
            finally:
                out_pdf.close()
                src.close()
        else:
            img = Image.open(p)
            pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                img, extension="pdf", lang=lang
            )
            out_path.write_bytes(pdf_bytes)
    except Exception as exc:
        print(f"[ocr] failed to write searchable PDF: {exc}")
        return None

    print(f"[ocr] wrote {out_path}")
    return out_path


def _build_ocr_pages(
    p: Path, *, lang: str, dpi: int
) -> list[
    tuple["Image.Image", list[tuple[int, int, int, int, str, int, int, int]]]
] | None:
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        print("[ocr] pytesseract/Pillow not available; skipping OCR pages")
        return None
    pages: list[
        tuple["Image.Image", list[tuple[int, int, int, int, str, int, int, int]]]
    ] = []

    def _auto_orient(img: "Image.Image") -> "Image.Image":
        if os.environ.get("OCR_DISABLE_OSD", "").strip() in {"1", "true", "yes"}:
            return img
        try:
            osd = pytesseract.image_to_osd(img)
        except Exception:
            return img
        rotate = 0
        for line in osd.splitlines():
            if line.startswith("Rotate:"):
                try:
                    rotate = int(line.split(":")[1].strip())
                except Exception:
                    rotate = 0
                break
        if rotate == 90:
            return img.rotate(90, expand=True)
        if rotate == 180:
            return img.rotate(180, expand=True)
        if rotate == 270:
            return img.rotate(270, expand=True)
        return img

    def _image_words(img: "Image.Image"):
        img = _auto_orient(img)
        try:
            psm = int(os.environ.get("OCR_PSM", "4"))
        except Exception:
            psm = 4
        data = pytesseract.image_to_data(
            img,
            lang=lang,
            output_type=pytesseract.Output.DICT,
            config=f"--psm {psm}",
        )
        words: list[tuple[int, int, int, int, str, int, int, int]] = []
        confs: list[float] = []
        n = len(data.get("text", []))
        for i in range(n):
            text = (data["text"][i] or "").strip()
            if not text:
                continue
            x, y, w, h = (
                int(data["left"][i]),
                int(data["top"][i]),
                int(data["width"][i]),
                int(data["height"][i]),
            )
            try:
                conf = float(data.get("conf", ["-1"])[i])
            except Exception:
                conf = -1.0
            if conf >= 0:
                confs.append(conf)
            block = int(data.get("block_num", [0])[i] or 0)
            par = int(data.get("par_num", [0])[i] or 0)
            line = int(data.get("line_num", [0])[i] or 0)
            words.append((x, y, x + w, y + h, text, block, par, line))
        avg_conf = sum(confs) / len(confs) if confs else 0.0
        return words, avg_conf

    try:
        if p.suffix.lower() in PDF_EXTS:
            src = fitz.open(str(p))
            try:
                zoom = dpi / 72.0
                mat = fitz.Matrix(zoom, zoom)
                for page in src:
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    img = Image.open(BytesIO(pix.tobytes("png")))
                    words, avg_conf = _image_words(img)
                    pages.append((img, words, avg_conf))
            finally:
                src.close()
        else:
            img = Image.open(p)
            words, avg_conf = _image_words(img)
            pages.append((img, words, avg_conf))
    except Exception as exc:
        print(f"[ocr] failed to generate OCR pages: {exc}")
        return None

    return pages


def _write_positioned_html(
    p: Path, *, pages: list[
        tuple["Image.Image", list[tuple[int, int, int, int, str, int, int, int]], float]
    ] | None
) -> Path | None:
    if not pages:
        return None

    out_path = p.with_suffix(".ocr.html")

    def _img_to_data_url(img: "Image.Image"):
        buf = BytesIO()
        img.save(buf, format="PNG")
        return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode("ascii")

    html_parts = [
        "<!doctype html>",
        "<html><head><meta charset=\"utf-8\">",
        "<style>",
        "body{margin:0;background:#111;color:#000;font-family:Helvetica,Arial,sans-serif;}",
        ".page{position:relative;margin:24px auto;box-shadow:0 8px 30px rgba(0,0,0,0.35);}",
        ".ocr-word{position:absolute;white-space:pre;color:#000;}",
        "</style></head><body>",
    ]

    for img, words, _avg_conf in pages:
        w, h = img.size
        data_url = _img_to_data_url(img)
        html_parts.append(
            f"<div class=\"page\" style=\"width:{w}px;height:{h}px;"
            f"background:url('{data_url}') no-repeat; background-size:{w}px {h}px;\">"
        )
        for x1, y1, x2, y2, text, _block, _par, _line in words:
            text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            font_size = max(8, y2 - y1)
            html_parts.append(
                "<span class=\"ocr-word\" "
                f"style=\"left:{x1}px;top:{y1}px;width:{x2 - x1}px;height:{y2 - y1}px;"
                f"font-size:{font_size}px;line-height:{y2 - y1}px;\">"
                f"{text}</span>"
            )
        html_parts.append("</div>")

    html_parts.append("</body></html>")
    out_path.write_text("\n".join(html_parts), encoding="utf-8")
    print(f"[ocr] wrote {out_path}")
    return out_path


def _write_layout_md(
    p: Path, *, pages: list[
        tuple["Image.Image", list[tuple[int, int, int, int, str, int, int, int]], float]
    ] | None
) -> str | None:
    if not pages:
        return None

    try:
        max_width = int(os.environ.get("OCR_LAYOUT_MD_WIDTH", "120"))
    except Exception:
        max_width = 120
    max_width = max(40, min(240, max_width))

    def _group_lines(
        words: list[tuple[int, int, int, int, str, int, int, int]]
    ) -> list[tuple[float, list[tuple[int, int, int, int, str]]]]:
        # Prefer Tesseract's line grouping for stable columns.
        buckets: dict[tuple[int, int, int], list[tuple[int, int, int, int, str]]] = {}
        ycs: dict[tuple[int, int, int], float] = {}
        for x1, y1, x2, y2, text, block, par, line in words:
            key = (block, par, line)
            buckets.setdefault(key, []).append((x1, y1, x2, y2, text))
            yc = (y1 + y2) / 2.0
            ycs[key] = (ycs.get(key, yc) + yc) / 2.0
        lines = [(ycs[k], buckets[k]) for k in buckets]
        return sorted(lines, key=lambda l: l[0])

    def _render_line(words: list[tuple[int, int, int, int, str]], page_w: int) -> str:
        words = sorted(words, key=lambda w: w[0])
        px_per_char = max(1.0, page_w / max(40, max_width))
        line = []
        cursor = 0
        first = True
        prev_x2 = 0
        for x1, _y1, x2, _y2, text in words:
            if not text.strip():
                continue
            if first:
                lead_spaces = int(round(x1 / px_per_char))
                line.append(" " * max(0, min(lead_spaces, max_width - 1)))
                cursor = len(line[0])
                first = False
            else:
                gap = max(0, x1 - prev_x2)
                spaces = int(round(gap / px_per_char))
                line.append(" " * max(1, min(spaces, max_width - cursor - 1)))
                cursor += len(line[-1])
            line.append(text)
            cursor += len(text)
            prev_x2 = x2
            if cursor >= max_width:
                break
        return "".join(line).rstrip()

    try:
        min_conf = float(os.environ.get("OCR_LAYOUT_MIN_CONF", "40"))
    except Exception:
        min_conf = 40.0

    blocks: list[str] = []
    for idx, (img, words, avg_conf) in enumerate(pages, start=1):
        w, _h = img.size
        blocks.append(f"--- Page {idx} ---")
        if avg_conf < min_conf:
            blocks.append(f"[layout skipped: low OCR confidence {avg_conf:.1f}]")
            continue
        lines = _group_lines(words)
        for _yc, line_words in lines:
            rendered = _render_line(line_words, w)
            if rendered:
                blocks.append(rendered)

    layout_text = "```text\n" + "\n".join(blocks) + "\n```\n"
    out_path = p.with_suffix(".ocr.layout.md")
    out_path.write_text(layout_text, encoding="utf-8")
    print(f"[ocr] wrote {out_path}")
    return layout_text


def _append_layout_to_md(p: Path, layout_text: str | None) -> None:
    if not layout_text:
        return
    out_path = p.with_suffix(".ocr.md")
    if not out_path.exists():
        return
    with out_path.open("a", encoding="utf-8") as f:
        f.write("\n\n")
        f.write("## OCR layout (approximate)\n\n")
        f.write(layout_text)


def _convert_to_pdf(input_path: Path, out_dir: Path) -> Path:
    soffice = shutil.which("soffice")
    if not soffice:
        raise RuntimeError("LibreOffice 'soffice' is required to convert Word files to PDF.")

    out_dir.mkdir(parents=True, exist_ok=True)
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(input_path),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(
            "LibreOffice conversion failed.\n"
            f"STDOUT:\n{result.stdout}\n\nSTDERR:\n{result.stderr}"
        )

    pdf_path = out_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise RuntimeError(f"LibreOffice did not produce PDF: {pdf_path}")
    return pdf_path
